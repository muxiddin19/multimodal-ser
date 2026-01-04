#!/usr/bin/env python3
"""
ACL 2026 Presentation: Sentimentogram
Enhanced version with figures, tables, and professional visual design
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

# Paths
BASE_DIR = "/home/muhiddin/ser"
FIGURES_DIR = os.path.join(BASE_DIR, "figures")
CAPTURE_DIR = os.path.join(BASE_DIR, "acl2026/capture")
OUTPUT_DIR = os.path.join(BASE_DIR, "presentation")

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme (Professional ACL-style)
COLORS = {
    'title': RGBColor(0, 51, 102),       # Dark blue
    'accent': RGBColor(0, 102, 153),     # Teal
    'text': RGBColor(51, 51, 51),        # Dark gray
    'highlight': RGBColor(204, 85, 0),   # Orange
    'success': RGBColor(34, 139, 34),    # Forest green
    'anger': RGBColor(220, 53, 69),      # Red
    'happy': RGBColor(255, 193, 7),      # Gold
    'sad': RGBColor(52, 152, 219),       # Blue
    'neutral': RGBColor(128, 128, 128),  # Gray
    'bg_light': RGBColor(245, 248, 250), # Light blue-gray
    'white': RGBColor(255, 255, 255),
    'black': RGBColor(0, 0, 0),
}

def add_background(slide, color=None):
    """Add a solid background to slide"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color or COLORS['bg_light']
    bg.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_header_bar(slide, title_text):
    """Add a colored header bar with title"""
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['title']
    header.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

def add_bullet_points(slide, items, left, top, width, height, font_size=20, color=None):
    """Add bullet points to slide"""
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = text_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or COLORS['text']
        p.space_after = Pt(8)

def add_notes(slide, notes_text):
    """Add presenter notes"""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

def add_table(slide, data, left, top, col_widths, row_height=0.4, header_color=None, font_size=14):
    """Add a formatted table"""
    rows = len(data)
    cols = len(data[0])

    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                   Inches(sum(col_widths)), Inches(rows * row_height)).table

    # Set column widths
    for i, width in enumerate(col_widths):
        table.columns[i].width = Inches(width)

    # Fill data
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)

            # Format
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = COLORS['text']
                paragraph.alignment = PP_ALIGN.CENTER

                # Header row styling
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = COLORS['white']

            # Header background
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color or COLORS['title']
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['white']

    return table

def add_icon_box(slide, text, left, top, width=2.5, height=1.5, color=None, icon_text=""):
    """Add an icon-style box with text"""
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = color or COLORS['accent']
    box.line.color.rgb = color or COLORS['accent']

    # Icon/number text
    if icon_text:
        icon_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.15), Inches(width), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    # Text
    text_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.6),
                                        Inches(width - 0.2), Inches(height - 0.7))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

# =============================================================================
# SLIDE 1: Title Slide
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['title'])

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Sentimentogram"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.333), Inches(1.2))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "Interpretable Multimodal Speech Emotion Recognition\nwith Human-Centered Design"
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(200, 220, 255)
p.alignment = PP_ALIGN.CENTER

# Conference
conf_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
tf = conf_box.text_frame
p = tf.paragraphs[0]
p.text = "ACL 2026"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLORS['happy']
p.alignment = PP_ALIGN.CENTER

# Author placeholder
author_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.5))
tf = author_box.text_frame
p = tf.paragraphs[0]
p.text = "Anonymous Authors"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(180, 200, 230)
p.alignment = PP_ALIGN.CENTER

add_notes(slide, """PRESENTER NOTES - Title Slide (30 seconds)

Welcome everyone. Today I'm presenting Sentimentogram - our work on interpretable multimodal speech emotion recognition.

Key opening points:
- This work focuses on HUMAN-CENTERED design, not just accuracy
- We address three critical gaps in current SER systems
- Our goal: Make emotion recognition UNDERSTANDABLE and USEFUL for real users

Transition: Let me start by explaining why this matters...""")

# =============================================================================
# SLIDE 2: Problem Statement
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "The Problem: Black-Box Emotion Recognition")

# Left side - Problem description
problem_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.5), Inches(4))
tf = problem_box.text_frame
tf.word_wrap = True

problems = [
    ("🔒 Opaque Decisions", "Models predict emotions but can't explain WHY"),
    ("📊 Raw Scores", "Probabilities are meaningless to end users"),
    ("👥 One-Size-Fits-All", "Cultural/personal preferences ignored"),
]

for i, (title, desc) in enumerate(problems):
    if i > 0:
        p = tf.add_paragraph()
        p.text = ""
        p.space_before = Pt(20)
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['anger']

    p = tf.add_paragraph()
    p.text = desc
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['text']

# Right side - visual representation
# Black box illustration
box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(2), Inches(3), Inches(2))
box.fill.solid()
box.fill.fore_color.rgb = COLORS['black']
box.line.color.rgb = COLORS['black']

box_text = slide.shapes.add_textbox(Inches(7), Inches(2.7), Inches(3), Inches(0.6))
tf = box_text.text_frame
p = tf.paragraphs[0]
p.text = "Black Box\nModel"
p.font.size = Pt(20)
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER

# Input arrow
arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.8), Inches(2.8), Inches(1.1), Inches(0.4))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = COLORS['accent']

# Output arrow
arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(10.1), Inches(2.8), Inches(1.1), Inches(0.4))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = COLORS['accent']

# Input label
in_label = slide.shapes.add_textbox(Inches(5.5), Inches(3.3), Inches(1.5), Inches(0.5))
tf = in_label.text_frame
p = tf.paragraphs[0]
p.text = "Audio +\nText"
p.font.size = Pt(12)
p.font.color.rgb = COLORS['text']
p.alignment = PP_ALIGN.CENTER

# Output label
out_label = slide.shapes.add_textbox(Inches(11), Inches(2.5), Inches(2), Inches(1))
tf = out_label.text_frame
p = tf.paragraphs[0]
p.text = "Angry: 0.73\nSad: 0.15\nHappy: 0.12"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']

# Question mark
q_box = slide.shapes.add_textbox(Inches(11.2), Inches(3.6), Inches(1), Inches(1))
tf = q_box.text_frame
p = tf.paragraphs[0]
p.text = "???"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = COLORS['anger']

# Bottom message
bottom = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12), Inches(0.5))
tf = bottom.text_frame
p = tf.paragraphs[0]
p.text = "Users need to UNDERSTAND and TRUST emotion predictions to use them effectively"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLORS['title']
p.alignment = PP_ALIGN.CENTER

add_notes(slide, """PRESENTER NOTES - Problem Statement (1 minute)

Current SER systems have three critical problems:

1. OPAQUE DECISIONS
   - Deep learning models are black boxes
   - Users can't understand why a prediction was made
   - Critical for clinical and accessibility applications

2. RAW PROBABILITY SCORES
   - Output like "Angry: 0.73" means nothing to end users
   - No visual representation that humans can quickly interpret
   - Subtitle systems still use plain text

3. NO PERSONALIZATION
   - Everyone gets the same visualization
   - Cultural differences ignored (red means anger in West, luck in East)
   - Individual preferences not considered

This matters because:
- Healthcare providers need explainable predictions
- Deaf/HoH users need emotion-rich subtitles
- Content creators want personalized tools

Transition: Our solution addresses all three...""")

# =============================================================================
# SLIDE 3: Our Solution - Three Pillars
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "Our Solution: Three Pillars of Human-Centered SER")

# Three pillars as colored boxes
pillars = [
    ("1", "Interpretable\nFusion", "Constrained gates that sum to 1:\n'76% audio, 24% text'", COLORS['accent']),
    ("2", "Emotion\nTypography", "Visual subtitles with\nemotion-specific styling", COLORS['success']),
    ("3", "Learned\nPersonalization", "Adapts to individual\npreferences, not stereotypes", COLORS['highlight']),
]

for i, (num, title, desc, color) in enumerate(pillars):
    left = 0.8 + i * 4.2

    # Pillar box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(left), Inches(1.8), Inches(3.8), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(left + 1.4), Inches(2.1), Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['white']
    circle.line.color.rgb = COLORS['white']

    num_text = slide.shapes.add_textbox(Inches(left + 1.4), Inches(2.2), Inches(1), Inches(0.9))
    tf = num_text.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    # Title
    title_text = slide.shapes.add_textbox(Inches(left + 0.2), Inches(3.3), Inches(3.4), Inches(1))
    tf = title_text.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Description
    desc_text = slide.shapes.add_textbox(Inches(left + 0.2), Inches(4.5), Inches(3.4), Inches(1.5))
    tf = desc_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(240, 240, 240)
    p.alignment = PP_ALIGN.CENTER

# Bottom arrow connecting all three
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(6.5), Inches(4.5), Inches(0.4))
arrow.fill.solid()
arrow.fill.fore_color.rgb = COLORS['title']

result_text = slide.shapes.add_textbox(Inches(9.2), Inches(6.35), Inches(3.5), Inches(0.7))
tf = result_text.text_frame
p = tf.paragraphs[0]
p.text = "Human-Centered SER"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLORS['title']

add_notes(slide, """PRESENTER NOTES - Three Pillars (1.5 minutes)

Our solution has THREE integrated pillars:

PILLAR 1: INTERPRETABLE FUSION
- Unlike black-box fusion, our gates are CONSTRAINED to sum to 1
- This enables statements like "76% audio, 24% text"
- Users can understand WHY a prediction was made
- Critical for clinical deployment

PILLAR 2: EMOTION TYPOGRAPHY
- Transform predictions into VISUAL subtitles
- Each emotion has distinct typography:
  * Anger: Bold, uppercase, red, shaking
  * Sadness: Italic, smaller, blue, fading
  * Happiness: Bouncy, gold, playful font
- Enables emotion perception without reading numbers

PILLAR 3: LEARNED PERSONALIZATION
- Rule-based cultural adaptation FAILS (43.8% vs 50.3% random!)
- We LEARN individual preferences from minimal feedback
- 10 comparisons = personalized visualization
- Avoids demographic stereotyping

KEY INSIGHT: These three pillars work TOGETHER:
- Interpretable fusion → informs typography decisions
- Typography → enables meaningful preference learning
- Personalization → makes typography effective for individuals

Transition: Let me show you how the interpretable fusion works...""")

# =============================================================================
# SLIDE 4: What Makes Us Different
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "Distinction from Prior Fusion Work")

# Two columns comparison
# Left: Prior Work
left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5))
tf = left_title.text_frame
p = tf.paragraphs[0]
p.text = "Prior Cross-Attention Work"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLORS['anger']

left_items = [
    "Focus: Maximize accuracy",
    "Fusion gates: Unconstrained (any values)",
    "Output: Probability scores only",
    "No human-interpretable explanations",
    "No visualization component",
    "No personalization"
]

for i, item in enumerate(left_items):
    item_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.1 + i * 0.6), Inches(5.3), Inches(0.5))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"✗ {item}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['text']

# Right: Our Work
right_title = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(0.5))
tf = right_title.text_frame
p = tf.paragraphs[0]
p.text = "Sentimentogram (Ours)"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLORS['success']

right_items = [
    "Focus: Interpretability + accuracy",
    "Fusion gates: Sum to 1 (percentages)",
    "Output: Visual emotion typography",
    "Per-sample modality attribution",
    "End-to-end visualization pipeline",
    "Learned user personalization"
]

for i, item in enumerate(right_items):
    item_box = slide.shapes.add_textbox(Inches(7.2), Inches(2.1 + i * 0.6), Inches(5.3), Inches(0.5))
    tf = item_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"✓ {item}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['text']

# Vertical divider
divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.4), Inches(1.5), Inches(0.05), Inches(4.2))
divider.fill.solid()
divider.fill.fore_color.rgb = COLORS['neutral']

# Key message at bottom
key_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(1), Inches(5.9), Inches(11.333), Inches(1))
key_box.fill.solid()
key_box.fill.fore_color.rgb = COLORS['title']

key_text = slide.shapes.add_textbox(Inches(1.2), Inches(6.05), Inches(11), Inches(0.8))
tf = key_text.text_frame
p = tf.paragraphs[0]
p.text = "Our novelty is NOT cross-attention itself, but the complete human-centered pipeline"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER

add_notes(slide, """PRESENTER NOTES - Distinction (1 minute)

This slide addresses a key question reviewers might have:
"Cross-attention is well-studied. What's new?"

CRITICAL DISTINCTION:
Prior work uses cross-attention to MAXIMIZE ACCURACY
We use constrained fusion for INTERPRETABILITY

Specific differences:
1. CONSTRAINED GATES
   - Prior: Gates can be any value (0.3, 0.8, 0.5 - doesn't help user)
   - Ours: Gates sum to 1 = percentage interpretation

2. COMPLETE PIPELINE
   - Prior: Model → probability → done
   - Ours: Model → interpretation → visualization → personalization

3. DESIGN PHILOSOPHY
   - Prior: "Higher accuracy is always better"
   - Ours: "Understandable predictions are more valuable"

KEY MESSAGE: The innovation is the PIPELINE, not the attention mechanism.

Transition: Let me show the architecture details...""")

# =============================================================================
# SLIDE 5: Architecture Overview
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "System Architecture")

# Feature extractors (left side)
# Text branch
text_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.5), Inches(1.8), Inches(2.2), Inches(1))
text_box.fill.solid()
text_box.fill.fore_color.rgb = COLORS['accent']
text_label = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(2.2), Inches(0.5))
tf = text_label.text_frame
p = tf.paragraphs[0]
p.text = "BERT\n(768-dim)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER

# Audio branch
audio_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.5), Inches(3.5), Inches(2.2), Inches(1))
audio_box.fill.solid()
audio_box.fill.fore_color.rgb = COLORS['highlight']
audio_label = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(2.2), Inches(0.5))
tf = audio_label.text_frame
p = tf.paragraphs[0]
p.text = "emotion2vec\n(1024-dim)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER

# Arrows from extractors
arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.8), Inches(2.1), Inches(0.8), Inches(0.3))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = COLORS['text']

arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.8), Inches(3.8), Inches(0.8), Inches(0.3))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = COLORS['text']

# VAD-Guided Cross-Attention (center)
vga_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(3.8), Inches(2.0), Inches(3), Inches(2.8))
vga_box.fill.solid()
vga_box.fill.fore_color.rgb = COLORS['success']
vga_label = slide.shapes.add_textbox(Inches(3.8), Inches(2.8), Inches(3), Inches(1.5))
tf = vga_label.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "VAD-Guided\nCross-Attention"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nPsychologically\ngrounded attention"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(220, 255, 220)
p.alignment = PP_ALIGN.CENTER

# Arrow to fusion
arrow3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.9), Inches(3.2), Inches(0.8), Inches(0.3))
arrow3.fill.solid()
arrow3.fill.fore_color.rgb = COLORS['text']

# Constrained Fusion (right-center)
fusion_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(7.8), Inches(2.0), Inches(2.5), Inches(2.8))
fusion_box.fill.solid()
fusion_box.fill.fore_color.rgb = COLORS['title']
fusion_label = slide.shapes.add_textbox(Inches(7.8), Inches(2.8), Inches(2.5), Inches(1.5))
tf = fusion_label.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Constrained\nAdaptive Fusion"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "\nαT + αA + αI = 1"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(200, 220, 255)
p.alignment = PP_ALIGN.CENTER

# Arrow to output
arrow4 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(10.4), Inches(3.2), Inches(0.8), Inches(0.3))
arrow4.fill.solid()
arrow4.fill.fore_color.rgb = COLORS['text']

# Output (right)
out_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(11.3), Inches(2.3), Inches(1.5), Inches(2.2))
out_box.fill.solid()
out_box.fill.fore_color.rgb = COLORS['anger']
out_label = slide.shapes.add_textbox(Inches(11.3), Inches(2.9), Inches(1.5), Inches(1))
tf = out_label.text_frame
p = tf.paragraphs[0]
p.text = "Emotion\nClass"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER

# Loss functions at bottom
loss_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(3), Inches(5.5), Inches(7.333), Inches(1.2))
loss_box.fill.solid()
loss_box.fill.fore_color.rgb = RGBColor(230, 230, 240)
loss_box.line.color.rgb = COLORS['title']

loss_label = slide.shapes.add_textbox(Inches(3.2), Inches(5.7), Inches(7), Inches(0.9))
tf = loss_label.text_frame
p = tf.paragraphs[0]
p.text = "ℒ = ℒfocal + λmicl·ℒMICL + λvad·ℒVAD"
p.font.size = Pt(22)
p.font.color.rgb = COLORS['title']
p.alignment = PP_ALIGN.CENTER
p = tf.add_paragraph()
p.text = "Focal Loss (class imbalance) + Hard Negative MICL + VAD Supervision"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']
p.alignment = PP_ALIGN.CENTER

add_notes(slide, """PRESENTER NOTES - Architecture (1.5 minutes)

Walk through the architecture left to right:

FEATURE EXTRACTION:
- Text: BERT-base-uncased (768-dim)
- Audio: emotion2vec-plus-large (1024-dim)
  * Pretrained on 40,000 hours of emotional speech
  * Already emotion-aware representations

VAD-GUIDED CROSS-ATTENTION:
- Standard bidirectional cross-attention
- BUT: Attention weights modulated by VAD predictions
- Valence-Arousal-Dominance provides psychological grounding
- High arousal emotions (anger) → focus on acoustic features
- This is NOT about accuracy, but INTERPRETABILITY

CONSTRAINED ADAPTIVE FUSION:
- Three gates: αT (text), αA (audio), αI (interaction)
- CRITICAL: αT + αA + αI = 1 (softmax constraint)
- This enables percentage interpretation
- Example: "This prediction used 76% audio, 24% text"

LOSS FUNCTION:
- Focal loss: Handles class imbalance
- Hard Negative MICL: Contrastive learning for confusion pairs
- VAD auxiliary: Weak supervision without ground truth

Transition: Let me explain the fusion constraint in detail...""")

# =============================================================================
# SLIDE 6: Constrained Adaptive Fusion
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "Constrained Adaptive Fusion: The Key to Interpretability")

# Left: Formula and explanation
formula_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(6), Inches(3))
tf = formula_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Fusion Equation:"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLORS['title']

p = tf.add_paragraph()
p.text = "\n  h = αT·hT + αA·hA + αI·(hT ⊙ hA)"
p.font.size = Pt(24)
p.font.color.rgb = COLORS['text']

p = tf.add_paragraph()
p.text = "\n\nConstraint:  αT + αA + αI = 1"
p.font.size = Pt(22)
p.font.color.rgb = COLORS['success']
p.font.bold = True

p = tf.add_paragraph()
p.text = "\n\nThis simple constraint enables:"
p.font.size = Pt(18)
p.font.color.rgb = COLORS['text']

benefits = [
    "Per-sample modality attribution",
    "\"76% audio, 24% text\" explanations",
    "Dataset-level analysis",
    "Debugging model behavior"
]
for benefit in benefits:
    p = tf.add_paragraph()
    p.text = f"  ✓ {benefit}"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['text']

# Right: Visual representation using stacked bars
# Bar chart visualization for fusion gates
bar_left = 7.5
bar_top = 2.5
bar_width = 4.5
bar_height = 0.8

# Text portion (24%)
text_bar_width = bar_width * 0.24
text_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(bar_left), Inches(bar_top),
                                   Inches(text_bar_width), Inches(bar_height))
text_bar.fill.solid()
text_bar.fill.fore_color.rgb = COLORS['accent']
text_bar.line.fill.background()

# Audio portion (76%)
audio_bar_width = bar_width * 0.76
audio_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(bar_left + text_bar_width), Inches(bar_top),
                                    Inches(audio_bar_width), Inches(bar_height))
audio_bar.fill.solid()
audio_bar.fill.fore_color.rgb = COLORS['highlight']
audio_bar.line.fill.background()

# Labels
label1 = slide.shapes.add_textbox(Inches(bar_left), Inches(bar_top + 0.9), Inches(text_bar_width), Inches(0.5))
tf = label1.text_frame
p = tf.paragraphs[0]
p.text = "Text: 24%"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['accent']
p.alignment = PP_ALIGN.CENTER

label2 = slide.shapes.add_textbox(Inches(bar_left + text_bar_width), Inches(bar_top + 0.9), Inches(audio_bar_width), Inches(0.5))
tf = label2.text_frame
p = tf.paragraphs[0]
p.text = "Audio: 76%"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLORS['highlight']
p.alignment = PP_ALIGN.CENTER

# Example box
example_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(7), Inches(5.5), Inches(5.5), Inches(1.3))
example_box.fill.solid()
example_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
example_box.line.color.rgb = COLORS['highlight']

example_text = slide.shapes.add_textbox(Inches(7.2), Inches(5.65), Inches(5.1), Inches(1.1))
tf = example_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Example: CREMA-D (acted emotions)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLORS['title']
p = tf.add_paragraph()
p.text = "Audio: 76.6% | Text: 23.1% | Interaction: 0.3%"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']

add_notes(slide, """PRESENTER NOTES - Constrained Fusion (1 minute)

This is THE key innovation for interpretability:

THE CONSTRAINT: αT + αA + αI = 1

Why this matters:
1. INTERPRETABLE AS PERCENTAGES
   - Without constraint: Gates could be 0.8, 0.6, 0.3 - meaningless to users
   - With constraint: "76% audio, 24% text" - instantly understandable

2. ENABLES EXPLANATION
   - User: "Why did you predict angry?"
   - System: "This prediction relied 76% on audio features"
   - Compare to black-box: "Because the model said so"

3. DATASET-LEVEL INSIGHTS
   - CREMA-D (acted): 76.6% audio - makes sense! Actors exaggerate vocal cues
   - IEMOCAP (natural): 54% text, 45% audio - balanced, natural speech is nuanced

4. NO ACCURACY SACRIFICE
   - Unconstrained: 92.21%
   - Constrained (ours): 93.02%
   - We IMPROVED accuracy while adding interpretability!

Transition: Now let's see the visualization component...""")

# =============================================================================
# SLIDE 7: Sentimentogram Typography (with image)
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "Emotion-Aware Typography: Sentimentogram")

# Try to add the image
img_path = os.path.join(CAPTURE_DIR, "sentimentogram_honest1.jpg")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(7))

# Typography rules on the right
rules_title = slide.shapes.add_textbox(Inches(7.8), Inches(1.5), Inches(5), Inches(0.5))
tf = rules_title.text_frame
p = tf.paragraphs[0]
p.text = "Typography Design Rules"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLORS['title']

# Emotion styles
emotions = [
    ("Anger", "Bebas Neue, BOLD, uppercase, 1.3× size, red", COLORS['anger']),
    ("Happiness", "Fredoka, playful, gold, bounce animation", COLORS['happy']),
    ("Sadness", "Merriweather, italic, smaller, blue, fade", COLORS['sad']),
    ("Neutral", "Poppins, clean, gray, standard size", COLORS['neutral']),
]

for i, (emotion, style, color) in enumerate(emotions):
    y = 2.2 + i * 0.9

    # Emotion label
    label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(7.8), Inches(y), Inches(1.8), Inches(0.6))
    label.fill.solid()
    label.fill.fore_color.rgb = color
    label_text = slide.shapes.add_textbox(Inches(7.8), Inches(y + 0.1), Inches(1.8), Inches(0.4))
    tf = label_text.text_frame
    p = tf.paragraphs[0]
    p.text = emotion
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Style description
    style_text = slide.shapes.add_textbox(Inches(9.8), Inches(y + 0.1), Inches(3), Inches(0.5))
    tf = style_text.text_frame
    p = tf.paragraphs[0]
    p.text = style
    p.font.size = Pt(13)
    p.font.color.rgb = COLORS['text']

# Key features at bottom
features_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.5), Inches(5.8), Inches(12.333), Inches(1.2))
features_box.fill.solid()
features_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
features_box.line.color.rgb = COLORS['accent']

features_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(12), Inches(1))
tf = features_text.text_frame
p = tf.paragraphs[0]
p.text = "✓ Real-time rendering   ✓ Character-level variation   ✓ Cultural adaptation   ✓ HTML output"
p.font.size = Pt(18)
p.font.color.rgb = COLORS['title']
p.alignment = PP_ALIGN.CENTER

add_notes(slide, """PRESENTER NOTES - Typography (1 minute)

Show the VISUAL demonstration of Sentimentogram:

IMAGE EXPLANATION:
- This is a real output from our system
- Each word is styled based on its predicted emotion
- Notice: "BEING HONEST" - bold, red, uppercase = anger
- Other words: gray, neutral styling
- The subtitle is READABLE while conveying emotion

TYPOGRAPHY DESIGN:
1. ANGER: Bebas Neue font, BOLD, UPPERCASE, 1.3× size, red
   - Visually "loud" and aggressive

2. HAPPINESS: Fredoka, playful rounded font, gold color
   - Bounce animation in video output

3. SADNESS: Merriweather serif, italic, smaller (0.92×), blue
   - Visually "deflated" and melancholic

4. NEUTRAL: Poppins clean sans-serif, gray
   - Doesn't distract from emotional words

KEY FEATURES:
- Real-time: 47ms latency (suitable for live video)
- Character-level: High-confidence words have per-character size variation
- Cultural: Western vs Eastern color symbolism profiles
- Output: Interactive HTML with synchronized video

Transition: But how do we personalize this for different users?...""")

# =============================================================================
# SLIDE 8: More Typography Examples
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "Typography Examples Across Emotions")

# Add multiple example images
example_images = [
    ("sentimentogram_balloons.jpg", 0.3, 1.5),
    ("sentimentogram_think.jpg", 6.6, 1.5),
    ("sentimentogram_gone.jpg", 0.3, 4.2),
]

for img_name, left, top in example_images:
    img_path = os.path.join(CAPTURE_DIR, img_name)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(6))

# Caption
caption = slide.shapes.add_textbox(Inches(6.6), Inches(4.5), Inches(6), Inches(2))
tf = caption.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Each example shows inline emotion styling:"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLORS['title']

p = tf.add_paragraph()
p.text = "\n• Emotion words highlighted WITHOUT boxes"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']

p = tf.add_paragraph()
p.text = "• Preserves natural reading flow"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']

p = tf.add_paragraph()
p.text = "• Emotion conveyed through typography alone"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']

add_notes(slide, """PRESENTER NOTES - More Examples (30 seconds)

These are additional real examples from our system:

Point out:
- Different emotions visible in each example
- Typography is INLINE - no boxes or labels needed
- Readability is maintained (users can still read the subtitle)
- Emotion is perceived ALONGSIDE the text content

This is what makes it "Sentimentogram":
- Like a cardiogram shows heart rhythm
- Sentimentogram shows emotional rhythm of speech

Transition: Now let's talk about personalization...""")

# =============================================================================
# SLIDE 9: Preference Learning
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "Why Rule-Based Personalization Fails")

# Left side: The problem
problem_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.5), Inches(0.5))
tf = problem_title.text_frame
p = tf.paragraphs[0]
p.text = "Rule-Based Approach"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLORS['anger']

# Rule examples
rules = [
    "IF region='East' THEN red=lucky",
    "IF age>50 THEN font_size=large",
    "IF gender='F' THEN colors=warm",
]

for i, rule in enumerate(rules):
    rule_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.3 + i * 0.6), Inches(5), Inches(0.5))
    tf = rule_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"• {rule}"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['text']

# Result box
result_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.7), Inches(4.3), Inches(4.8), Inches(1.2))
result_box.fill.solid()
result_box.fill.fore_color.rgb = RGBColor(255, 230, 230)
result_box.line.color.rgb = COLORS['anger']

result_text = slide.shapes.add_textbox(Inches(0.9), Inches(4.5), Inches(4.4), Inches(0.9))
tf = result_text.text_frame
p = tf.paragraphs[0]
p.text = "Result: 43.8% accuracy"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLORS['anger']
p = tf.add_paragraph()
p.text = "WORSE than random! (p=0.014)"
p.font.size = Pt(16)
p.font.color.rgb = COLORS['text']

# Right side: Our approach
solution_title = slide.shapes.add_textbox(Inches(6.5), Inches(1.6), Inches(6), Inches(0.5))
tf = solution_title.text_frame
p = tf.paragraphs[0]
p.text = "Learned Personalization"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLORS['success']

# Process
process = [
    "1. Show 10-15 pairwise comparisons",
    "2. User picks preferred style",
    "3. Learn preference function (N=50)",
    "4. Personalize for new content",
]

for i, step in enumerate(process):
    step_box = slide.shapes.add_textbox(Inches(6.7), Inches(2.3 + i * 0.6), Inches(5.5), Inches(0.5))
    tf = step_box.text_frame
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['text']

# Result box
result_box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(6.7), Inches(4.3), Inches(4.8), Inches(1.2))
result_box2.fill.solid()
result_box2.fill.fore_color.rgb = RGBColor(230, 255, 230)
result_box2.line.color.rgb = COLORS['success']

result_text2 = slide.shapes.add_textbox(Inches(6.9), Inches(4.5), Inches(4.4), Inches(0.9))
tf = result_text2.text_frame
p = tf.paragraphs[0]
p.text = "Result: 61.2% accuracy"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLORS['success']
p = tf.add_paragraph()
p.text = "vs Bradley-Terry 52.8% (p<0.001)"
p.font.size = Pt(16)
p.font.color.rgb = COLORS['text']

# Bottom insight
insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.5), Inches(5.9), Inches(12.333), Inches(1))
insight_box.fill.solid()
insight_box.fill.fore_color.rgb = COLORS['title']

insight_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.1), Inches(12), Inches(0.7))
tf = insight_text.text_frame
p = tf.paragraphs[0]
p.text = "Insight: Learn from individuals, don't stereotype from demographics"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER

add_notes(slide, """PRESENTER NOTES - Preference Learning (1.5 minutes)

This is a SURPRISING and IMPORTANT finding:

RULE-BASED PERSONALIZATION FAILS:
- Common approach: Use demographics to infer preferences
- "Eastern users prefer red for positive emotions"
- "Older users need larger fonts"
- Result: 43.8% accuracy - WORSE THAN RANDOM (p=0.014)!

Why does it fail?
- Demographics don't predict individual preferences
- Stereotyping leads to wrong assumptions
- Individual variation is larger than group patterns

OUR LEARNED APPROACH (N=50 real users, 1500 comparisons):
1. Show user 10-12 pairwise comparisons (~3 minutes)
2. User simply picks which subtitle style they prefer
3. Train lightweight logistic regression model
4. Personalize all future content

Stronger baselines tested:
- Bradley-Terry hierarchical: 52.8%
- Collaborative filtering: 53.5%
- Contextual logistic: 52.1%
- Our learned approach: 61.2% (p<0.001)

Direct A/B study (N=15 held-out users):
- Satisfaction: +8.7% (p=0.001)
- Comprehension: +5.8% (p<0.001)

BROADER IMPLICATION:
- This finding applies beyond SER
- NLP interfaces should learn from user feedback
- Not rely on demographic assumptions
- Especially important for accessibility applications

Transition: Now let's see the quantitative results...""")

# =============================================================================
# SLIDE 10: Main Results Table
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "Main Results: Cross-Dataset Performance")

# Results table
results_data = [
    ["Method", "IEMOCAP-4", "IEMOCAP-5", "IEMOCAP-6", "CREMA-D", "MELD"],
    ["BERT-only (Text)", "63.67", "52.87", "47.72", "28.96", "56.47"],
    ["emotion2vec (Audio)", "91.27", "76.22", "65.65", "91.84", "52.94"],
    ["Concatenation", "90.74", "76.51", "68.91", "92.09", "62.91"],
    ["Standard Cross-Attn", "89.33", "73.76", "66.14", "91.99", "63.10"],
    ["Adaptive (Unconstrained)", "92.21", "75.66", "65.97", "92.09", "59.97"],
    ["Sentimentogram (Ours)", "93.02", "77.97", "68.75", "92.90", "63.66"],
]

table = add_table(slide, results_data, 0.5, 1.6, [3, 1.8, 1.8, 1.8, 1.8, 1.8],
                  row_height=0.55, font_size=14)

# Highlight our results row
for col_idx in range(6):
    cell = table.cell(6, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(230, 255, 230)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLORS['success']

# Key findings
findings_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.5))
tf = findings_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Key Findings:"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLORS['title']

p = tf.add_paragraph()
p.text = "• Best or competitive on ALL datasets and configurations"
p.font.size = Pt(16)
p.font.color.rgb = COLORS['text']

p = tf.add_paragraph()
p.text = "• 93.02% UA on IEMOCAP-4: Interpretable fusion does NOT sacrifice accuracy"
p.font.size = Pt(16)
p.font.color.rgb = COLORS['text']

p = tf.add_paragraph()
p.text = "• Generalizes across: Acted (CREMA-D) | Natural (IEMOCAP) | TV dialogue (MELD)"
p.font.size = Pt(16)
p.font.color.rgb = COLORS['text']

add_notes(slide, """PRESENTER NOTES - Main Results (1 minute)

Walk through the table:

BASELINES:
- BERT-only: Text is weak alone for emotion
- emotion2vec: Strong audio baseline (91.27% on IEMOCAP-4)
- Concatenation: Simple fusion works reasonably well
- Standard cross-attention: Actually WORSE than concatenation!
- Unconstrained adaptive: 92.21% - our accuracy competitor

OUR RESULTS (bottom row, highlighted):
- IEMOCAP-4: 93.02% - BEST, +0.81% over unconstrained
- IEMOCAP-5: 77.97% - BEST, +1.46% improvement
- CREMA-D: 92.90% - BEST
- MELD: 63.66% - BEST

KEY MESSAGE:
"We achieve BETTER accuracy with INTERPRETABLE fusion"
- Sum-to-one constraint adds interpretability
- Does NOT hurt performance
- Actually HELPS in most cases

GENERALIZATION:
- IEMOCAP: Natural dyadic conversations
- CREMA-D: Acted/scripted emotions
- MELD: Multi-party TV dialogue
- All different domains, consistent improvement

Transition: Let's look at ASR robustness...""")

# =============================================================================
# SLIDE 10b: ASR Robustness (NEW)
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "Robustness to ASR Errors")

# ASR Results table
asr_data = [
    ["Condition", "WER", "UA", "Δ UA"],
    ["Ground Truth", "0%", "76.26%", "—"],
    ["Whisper-large-v3", "44.3%", "75.30%", "-0.96%"],
]

add_table(slide, asr_data, 1.5, 1.8, [3.5, 2, 2, 2], row_height=0.55, font_size=16)

# Key finding box
finding_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(1.5), Inches(3.5), Inches(10), Inches(1.5))
finding_box.fill.solid()
finding_box.fill.fore_color.rgb = COLORS['success']

finding_text = slide.shapes.add_textbox(Inches(1.7), Inches(3.7), Inches(9.6), Inches(1.2))
tf = finding_text.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Key Finding: Only -0.96% UA drop with 44% WER"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLORS['white']

p = tf.add_paragraph()
p.text = "→ Model is robust to real-world ASR errors"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(230, 255, 230)

# Why robust box
why_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.5), Inches(5.3), Inches(12.333), Inches(1.6))
why_box.fill.solid()
why_box.fill.fore_color.rgb = COLORS['white']
why_box.line.color.rgb = COLORS['title']

why_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(12), Inches(1.3))
tf = why_text.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Why Robust?"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLORS['title']

p = tf.add_paragraph()
p.text = "• Multimodal fusion compensates: Audio features remain unaffected"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']

p = tf.add_paragraph()
p.text = "• BERT tolerates noise: Semantic similarity preserved despite word errors"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']

p = tf.add_paragraph()
p.text = "• Tested with real Whisper transcriptions, not simulated noise"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']

add_notes(slide, """PRESENTER NOTES - ASR Robustness (45 seconds)

CRITICAL for real-world deployment:

EXPERIMENTAL SETUP:
- Used Whisper-large-v3 to transcribe IEMOCAP audio
- Compared ground-truth text vs ASR transcripts
- Real WER: 44.3% (quite noisy!)

RESULTS:
- Ground truth: 76.26% UA
- With ASR errors: 75.30% UA
- Degradation: Only 0.96%!

WHY SO ROBUST?
1. Multimodal fusion: Audio features unaffected by text errors
2. BERT is robust: Semantic meaning preserved despite word changes
3. Learned redundancy: Model learns to use both modalities

IMPORTANCE:
- Real-world audio has ASR errors
- If model only worked with clean text, it would be useless
- This shows practical deployability

Transition: Now let's see the ablation study...""")

# =============================================================================
# SLIDE 11: Ablation Study
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "Ablation Study: Component Contributions")

# Two-column layout
# Left: Accuracy impact
acc_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.5), Inches(0.5))
tf = acc_title.text_frame
p = tf.paragraphs[0]
p.text = "Accuracy Impact"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLORS['title']

acc_items = [
    "Full model: 93.02%",
    "- VAD auxiliary loss: 91.82% (-1.2%)",
    "Audio-only: 91.27% (p=0.02)",
    "Text-only: 63.67% (p<0.001)",
]

acc_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.5), Inches(2))
tf = acc_box.text_frame
for i, item in enumerate(acc_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = item
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['text']
    p.space_after = Pt(8)

# Right: Interpretability value
int_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(6), Inches(0.5))
tf = int_title.text_frame
p = tf.paragraphs[0]
p.text = "Interpretability Value (Primary)"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = COLORS['success']

int_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.2), Inches(5.5), Inches(2.5))
tf = int_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Unconstrained fusion:"
p.font.size = Pt(16)
p.font.color.rgb = COLORS['text']
p = tf.add_paragraph()
p.text = "  92.21% accuracy, NO interpretability"
p.font.size = Pt(16)
p.font.color.rgb = COLORS['text']

p = tf.add_paragraph()
p.text = "\nConstrained fusion (ours):"
p.font.size = Pt(16)
p.font.color.rgb = COLORS['text']
p = tf.add_paragraph()
p.text = "  93.02% accuracy + interpretability"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLORS['success']

p = tf.add_paragraph()
p.text = '\n→ "76% audio, 24% text" explanations'
p.font.size = Pt(16)
p.font.color.rgb = COLORS['text']

# Synergy explanation at bottom
synergy_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.5), Inches(4.8), Inches(12.333), Inches(2))
synergy_box.fill.solid()
synergy_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
synergy_box.line.color.rgb = COLORS['highlight']

synergy_text = slide.shapes.add_textbox(Inches(0.7), Inches(5), Inches(12), Inches(1.7))
tf = synergy_text.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Why Component Synergy is Expected:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLORS['highlight']

p = tf.add_paragraph()
p.text = "\nComponents were designed to work together: VAD guides attention → Attention informs fusion → Fusion enables MICL"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']

p = tf.add_paragraph()
p.text = "The whole exceeding the sum of parts indicates successful integration, not a limitation."
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']

add_notes(slide, """PRESENTER NOTES - Ablation Study (1 minute)

TWO AXES of evaluation:

ACCURACY IMPACT:
- Multimodal fusion is essential (audio-only and text-only much worse)
- VAD auxiliary loss provides useful regularization (-1.2% without it)
- Individual components show synergistic effects

INTERPRETABILITY VALUE (more important!):
- Unconstrained fusion: 92.21% accuracy, but gates are meaningless
- Constrained fusion: 93.02% accuracy + interpretability
- We get BOTH better accuracy AND explanation capability

ADDRESSING POTENTIAL CRITICISM:
"But individual components don't show isolated effects..."

Response: This is EXPECTED and INTENTIONAL
- Components were designed to complement each other
- VAD → Attention → Fusion → MICL is a pipeline
- The whole should exceed the sum of parts
- This indicates successful INTEGRATION

KEY MESSAGE:
"Interpretability is the PRIMARY contribution, not just accuracy"

Transition: Let me show the fusion gate analysis...""")

# =============================================================================
# SLIDE 12: Fusion Gate Analysis
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "Fusion Gate Analysis: Dataset Insights")

# Fusion gate table
fusion_data = [
    ["Dataset", "Text", "Audio", "Interaction"],
    ["IEMOCAP 5-class", "54.3%", "45.5%", "0.2%"],
    ["IEMOCAP 6-class", "41.4%", "58.4%", "0.2%"],
    ["CREMA-D", "23.1%", "76.6%", "0.3%"],
]

add_table(slide, fusion_data, 0.5, 1.6, [3.5, 2, 2, 2], row_height=0.55, font_size=16)

# Visual bar representation
bar_top = 3.8
datasets_vis = [
    ("IEMOCAP 5-class", 54.3, 45.5, 0.2),
    ("IEMOCAP 6-class", 41.4, 58.4, 0.2),
    ("CREMA-D", 23.1, 76.6, 0.3),
]

for i, (name, text_pct, audio_pct, int_pct) in enumerate(datasets_vis):
    y = bar_top + i * 0.9

    # Label
    label = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(2.5), Inches(0.4))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text']

    # Text bar
    bar_width = 8.5 * (text_pct / 100)
    if bar_width > 0:
        text_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Inches(3), Inches(y), Inches(bar_width), Inches(0.35))
        text_bar.fill.solid()
        text_bar.fill.fore_color.rgb = COLORS['accent']
        text_bar.line.fill.background()

    # Audio bar
    audio_bar_width = 8.5 * (audio_pct / 100)
    if audio_bar_width > 0:
        audio_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Inches(3 + bar_width), Inches(y),
                                            Inches(audio_bar_width), Inches(0.35))
        audio_bar.fill.solid()
        audio_bar.fill.fore_color.rgb = COLORS['highlight']
        audio_bar.line.fill.background()

# Legend
legend_text = slide.shapes.add_textbox(Inches(3), Inches(6.5), Inches(8), Inches(0.4))
tf = legend_text.text_frame
p = tf.paragraphs[0]
p.text = "■ Text        ■ Audio"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']

# Insight box
insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(7.5), Inches(4.5), Inches(5.3), Inches(2))
insight_box.fill.solid()
insight_box.fill.fore_color.rgb = COLORS['title']

insight_text = slide.shapes.add_textbox(Inches(7.7), Inches(4.7), Inches(4.9), Inches(1.7))
tf = insight_text.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Key Insight"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLORS['white']

p = tf.add_paragraph()
p.text = "\nCREMA-D (acted): 76.6% audio"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(200, 220, 255)

p = tf.add_paragraph()
p.text = "→ Actors exaggerate vocal cues"
p.font.size = Pt(13)
p.font.color.rgb = RGBColor(180, 200, 230)

p = tf.add_paragraph()
p.text = "\nIEMOCAP (natural): ~50/50"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(200, 220, 255)

p = tf.add_paragraph()
p.text = "→ Natural speech needs both modalities"
p.font.size = Pt(13)
p.font.color.rgb = RGBColor(180, 200, 230)

add_notes(slide, """PRESENTER NOTES - Fusion Analysis (1 minute)

This is WHERE interpretability becomes valuable:

CREMA-D (Acted Speech):
- 76.6% audio, 23.1% text
- Makes sense! Actors EXAGGERATE vocal expressions
- Emotion is mostly in HOW they speak, not WHAT they say

IEMOCAP (Natural Conversation):
- 54.3% text, 45.5% audio (5-class)
- More balanced fusion
- Natural speech requires understanding BOTH content AND delivery

INTERACTION GATES:
- Minimal (<1%) across all datasets
- Suggests additive rather than multiplicative integration
- Text and audio provide complementary, not interactive, signals

WHY THIS MATTERS:
- Without interpretability, these insights are invisible
- With constrained gates, we can EXPLAIN model behavior
- Useful for debugging, domain adaptation, clinical deployment

Example explanation:
"The model predicted ANGRY because it relied 76% on vocal features,
which showed elevated pitch and fast speaking rate."

Transition: Let's see the user study results...""")

# =============================================================================
# SLIDE 13: User Study Results
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "User Study: Typography Evaluation (N=30)")

# Key metrics boxes
metrics = [
    ("98%", "Reading Speed", "Preserved", COLORS['success']),
    ("84.2%", "Emotion Recognition", "vs 61.3% baseline", COLORS['accent']),
    ("87.3%", "Anger from Typography", "Alone (no audio)", COLORS['anger']),
    ("79.2%", "Happy from Typography", "Alone (no audio)", COLORS['happy']),
]

for i, (value, label, sublabel, color) in enumerate(metrics):
    left = 0.5 + i * 3.2

    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(left), Inches(1.7), Inches(3), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = color

    # Value
    val_text = slide.shapes.add_textbox(Inches(left), Inches(1.9), Inches(3), Inches(0.8))
    tf = val_text.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Label
    lab_text = slide.shapes.add_textbox(Inches(left), Inches(2.7), Inches(3), Inches(0.5))
    tf = lab_text.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Sublabel
    sub_text = slide.shapes.add_textbox(Inches(left), Inches(3.2), Inches(3), Inches(0.5))
    tf = sub_text.text_frame
    p = tf.paragraphs[0]
    p.text = sublabel
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(230, 230, 230)
    p.alignment = PP_ALIGN.CENTER

# Study details
details_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.5), Inches(4.3), Inches(12.333), Inches(2.5))
details_box.fill.solid()
details_box.fill.fore_color.rgb = COLORS['white']
details_box.line.color.rgb = COLORS['title']

details_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.5), Inches(12), Inches(2.2))
tf = details_text.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Study Design:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLORS['title']

p = tf.add_paragraph()
p.text = "• Within-subjects design with 30 participants"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']

p = tf.add_paragraph()
p.text = "• Compared: Plain subtitles vs. Full typography vs. Reduced typography"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']

p = tf.add_paragraph()
p.text = "• Measured: Reading speed, emotion recognition accuracy, discriminability"
p.font.size = Pt(15)
p.font.color.rgb = COLORS['text']

p = tf.add_paragraph()
p.text = "\nKey Finding: Typography enables emotion perception from text alone (87.3% anger, 79.2% happy)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLORS['success']

p = tf.add_paragraph()
p.text = "All results statistically significant (p < 0.001)"
p.font.size = Pt(14)
p.font.color.rgb = COLORS['text']

add_notes(slide, """PRESENTER NOTES - User Study (1 minute)

KEY FINDINGS from our user study:

1. READING SPEED PRESERVED (98%)
   - Typography doesn't slow reading
   - Critical for practical deployment
   - Users can read at normal pace

2. EMOTION RECOGNITION IMPROVED (84.2% vs 61.3%)
   - +22.9% improvement with typography
   - Users correctly identify more emotions
   - Typography adds information value

3. TYPOGRAPHY-ONLY RECOGNITION:
   - Anger: 87.3% accuracy from typography alone
   - Happiness: 79.2% accuracy
   - Users can identify emotion WITHOUT audio
   - Critical for deaf/HoH accessibility

STATISTICAL SIGNIFICANCE:
- All comparisons p < 0.001
- Bootstrap confidence intervals
- Within-subjects design reduces variance

STUDY DESIGN:
- N = 30 participants
- Balanced demographics
- Three conditions: Plain, Full, Reduced
- Randomized order

Transition: Let's discuss system latency...""")

# =============================================================================
# SLIDE 14: System Latency
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "Real-Time Capability: System Latency")

# Latency breakdown table
latency_data = [
    ["Component", "Latency", "Notes"],
    ["Audio feature extraction", "12ms", "emotion2vec inference"],
    ["Text feature extraction", "8ms", "BERT inference"],
    ["Model forward pass", "15ms", "Cross-attention + fusion"],
    ["Typography rendering", "12ms", "HTML generation"],
    ["Total", "47ms", "< 50ms threshold"],
]

add_table(slide, latency_data, 0.5, 1.6, [4, 2, 5.5], row_height=0.55, font_size=15)

# Highlight total row
# (Note: This would need additional formatting in real implementation)

# Real-time visualization
rt_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.5), Inches(4.5), Inches(5.5), Inches(2.3))
rt_box.fill.solid()
rt_box.fill.fore_color.rgb = COLORS['success']

rt_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.7), Inches(5.1), Inches(2))
tf = rt_text.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Real-Time Ready"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLORS['white']

p = tf.add_paragraph()
p.text = "\n✓ 47ms total latency"
p.font.size = Pt(18)
p.font.color.rgb = COLORS['white']

p = tf.add_paragraph()
p.text = "✓ Suitable for live video"
p.font.size = Pt(18)
p.font.color.rgb = COLORS['white']

p = tf.add_paragraph()
p.text = "✓ Streaming applications"
p.font.size = Pt(18)
p.font.color.rgb = COLORS['white']

# Applications box
app_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(6.5), Inches(4.5), Inches(6.333), Inches(2.3))
app_box.fill.solid()
app_box.fill.fore_color.rgb = COLORS['accent']

app_text = slide.shapes.add_textbox(Inches(6.7), Inches(4.7), Inches(5.9), Inches(2))
tf = app_text.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Practical Applications"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLORS['white']

apps = ["Live streaming subtitles", "Video conferencing", "Real-time feedback tools", "Accessibility systems"]
for app in apps:
    p = tf.add_paragraph()
    p.text = f"• {app}"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(220, 240, 255)

add_notes(slide, """PRESENTER NOTES - System Latency (30 seconds)

LATENCY BREAKDOWN:
- Audio extraction: 12ms (emotion2vec is efficient)
- Text extraction: 8ms (BERT is fast for short sequences)
- Model forward: 15ms (attention + fusion)
- Typography: 12ms (HTML rendering)
- TOTAL: 47ms

WHY 47ms MATTERS:
- Human perception threshold: ~100ms
- 47ms is well below perceptible delay
- Suitable for real-time applications

PRACTICAL DEPLOYMENT:
- Live streaming: Emotion-styled subtitles
- Video calls: Real-time emotion feedback
- Therapy: Immediate emotional reflection
- Accessibility: Live captions with emotion

Hardware: All measurements on single NVIDIA RTX 3090
Could be faster with optimization/quantization

Transition: Let me discuss limitations honestly...""")

# =============================================================================
# SLIDE 15: Limitations
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "Limitations and Future Work")

# Limitations
limitations = [
    ("No Video Modality", "Facial expressions provide valuable cues", "Future: Add visual modality"),
    ("English Only", "Cross-lingual generalization untested", "Future: Multilingual evaluation"),
    ("Utterance-Level", "No conversational context modeling", "Future: Dialogue history"),
    ("Synergistic Components", "Hard to isolate individual effects", "By design: Integrated pipeline"),
]

for i, (title, desc, future) in enumerate(limitations):
    y = 1.6 + i * 1.35

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(3.5), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"• {title}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['anger']

    # Description
    desc_box = slide.shapes.add_textbox(Inches(4.2), Inches(y), Inches(4.5), Inches(0.4))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['text']

    # Future
    future_box = slide.shapes.add_textbox(Inches(9), Inches(y), Inches(4), Inches(0.4))
    tf = future_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"→ {future}"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['success']

# Honest acknowledgment box
honest_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.5), Inches(5.9), Inches(12.333), Inches(1.2))
honest_box.fill.solid()
honest_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
honest_box.line.color.rgb = COLORS['highlight']

honest_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.1), Inches(12), Inches(0.9))
tf = honest_text.text_frame
p = tf.paragraphs[0]
p.text = "We acknowledge these limitations openly. Our primary contribution is the human-centered"
p.font.size = Pt(16)
p.font.color.rgb = COLORS['text']
p = tf.add_paragraph()
p.text = "design philosophy, which we believe is more impactful than incremental accuracy gains."
p.font.size = Pt(16)
p.font.color.rgb = COLORS['text']

add_notes(slide, """PRESENTER NOTES - Limitations (1 minute)

Be HONEST about limitations:

1. NO VIDEO MODALITY
   - Facial expressions are valuable for emotion
   - We focused on audio+text for this work
   - Future: Add visual processing

2. ENGLISH ONLY
   - All datasets are English
   - Cross-lingual transfer untested
   - Future: Multilingual evaluation (follow UniSER)

3. UTTERANCE-LEVEL
   - We don't model dialogue history
   - Conversation context helps ambiguous cases
   - Future: Add dialogue modeling

4. SYNERGISTIC COMPONENTS
   - Components work together, hard to isolate
   - This is BY DESIGN, not a bug
   - Integrated pipeline is the point

KEY MESSAGE:
"We're honest about what we don't do.
Our contribution is the PHILOSOPHY of human-centered design,
not just incremental accuracy improvements."

Transition: Let me summarize...""")

# =============================================================================
# SLIDE 16: Conclusion
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_header_bar(slide, "Conclusion")

# Summary boxes
summaries = [
    ("Interpretable Fusion", "Constrained gates sum to 1\n→ '76% audio, 24% text' explanations", COLORS['accent']),
    ("Emotion Typography", "Visual subtitles with emotion-specific styling\n→ 84.2% recognition improvement", COLORS['success']),
    ("Learned Personalization", "Rule-based fails (43.8%)\n→ Learned: 61.2% (N=50, p<0.001)", COLORS['highlight']),
]

for i, (title, desc, color) in enumerate(summaries):
    left = 0.5 + i * 4.2

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(left), Inches(1.6), Inches(4), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = color

    title_text = slide.shapes.add_textbox(Inches(left + 0.1), Inches(1.75), Inches(3.8), Inches(0.5))
    tf = title_text.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    desc_text = slide.shapes.add_textbox(Inches(left + 0.1), Inches(2.3), Inches(3.8), Inches(1.4))
    tf = desc_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(240, 240, 240)
    p.alignment = PP_ALIGN.CENTER

# Key numbers
numbers = [
    ("93.02%", "IEMOCAP-4 UA"),
    ("77.97%", "IEMOCAP-5 UA"),
    ("92.90%", "CREMA-D UA"),
    ("47ms", "Latency"),
]

for i, (num, label) in enumerate(numbers):
    left = 0.8 + i * 3.2

    num_text = slide.shapes.add_textbox(Inches(left), Inches(4.1), Inches(2.8), Inches(0.6))
    tf = num_text.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['title']
    p.alignment = PP_ALIGN.CENTER

    label_text = slide.shapes.add_textbox(Inches(left), Inches(4.6), Inches(2.8), Inches(0.4))
    tf = label_text.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['text']
    p.alignment = PP_ALIGN.CENTER

# Takeaway
takeaway_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.5), Inches(5.4), Inches(12.333), Inches(1.5))
takeaway_box.fill.solid()
takeaway_box.fill.fore_color.rgb = COLORS['title']

takeaway_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.6), Inches(12), Inches(1.2))
tf = takeaway_text.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Key Takeaway:"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLORS['happy']

p = tf.add_paragraph()
p.text = "Human-centered design doesn't sacrifice accuracy—it enhances it while making"
p.font.size = Pt(18)
p.font.color.rgb = COLORS['white']

p = tf.add_paragraph()
p.text = "emotion recognition understandable and useful for real users."
p.font.size = Pt(18)
p.font.color.rgb = COLORS['white']

add_notes(slide, """PRESENTER NOTES - Conclusion (1 minute)

SUMMARIZE the three pillars:

1. INTERPRETABLE FUSION
   - Constrained gates enable percentage explanations
   - "76% audio, 24% text" is understandable
   - Critical for clinical/accessibility deployment

2. EMOTION TYPOGRAPHY
   - Visual subtitles convey emotion through styling
   - +22.9% emotion recognition improvement
   - Enables perception without audio (accessibility)

3. LEARNED PERSONALIZATION (N=50 real users, 1500 comparisons)
   - Rule-based cultural adaptation FAILS (43.8%, p=0.014)
   - Learned approach: 61.2% vs Bradley-Terry 52.8% (p<0.001)
   - A/B study: +8.7% satisfaction, +5.8% comprehension
   - Learn from individuals, not stereotypes

KEY NUMBERS:
- 93.02% on IEMOCAP-4 (competitive with SOTA)
- 77.97% on IEMOCAP-5
- 47ms latency (real-time capable)
- ASR robust: Only -0.96% UA with 44% WER

FINAL MESSAGE:
"We showed that human-centered design is NOT a tradeoff—
you can have interpretability, visualization, personalization,
AND competitive accuracy. This is the future of NLP systems."

Transition: Thank you, questions?""")

# =============================================================================
# SLIDE 17: Thank You / Q&A
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, COLORS['title'])

# Thank you
thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
tf = thanks_box.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = COLORS['white']
p.alignment = PP_ALIGN.CENTER

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.8))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "Questions?"
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(200, 220, 255)
p.alignment = PP_ALIGN.CENTER

# Resources
resources = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1.2))
tf = resources.text_frame
p = tf.paragraphs[0]
p.text = "Code & Data: github.com/[anonymous]/sentimentogram"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(180, 200, 230)
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "Paper: ACL 2026 Proceedings"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(180, 200, 230)
p.alignment = PP_ALIGN.CENTER

add_notes(slide, """PRESENTER NOTES - Thank You (open for Q&A)

ANTICIPATED QUESTIONS:

Q: "How does this compare to GPT-4 or LLM-based emotion recognition?"
A: LLMs are powerful but lack interpretability and real-time capability.
   Our approach is complementary—could use LLM for text understanding,
   but fusion interpretability remains our unique contribution.

Q: "Why not use video modality?"
A: Focused scope for this work. Video adds complexity.
   Audio+text covers many use cases. Future work will add video.

Q: "How do you handle mixed emotions?"
A: Current model predicts dominant emotion.
   Future: Multi-label prediction or VAD regression.

Q: "Is 30 participants enough for user study?"
A: Within-subjects design reduces variance.
   Effect sizes are large (p < 0.001).
   Standard for HCI/NLP user studies.

Q: "Real deployment considerations?"
A: 47ms latency is practical. Privacy concerns addressed.
   Would need user consent for emotion tracking.

Thank the audience and invite specific questions!""")

# =============================================================================
# Save presentation
# =============================================================================
output_path = os.path.join(OUTPUT_DIR, "ACL2026_Sentimentogram.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
